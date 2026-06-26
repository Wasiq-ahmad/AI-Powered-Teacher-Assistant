import io
from typing import Optional
import PyPDF2
import docx
import openpyxl
from pptx import Presentation

class DocumentExtractor:
    @staticmethod
    def extract_text(file_content: bytes, filename: str) -> str:
        ext = filename.split('.')[-1].lower()
        
        try:
            if ext == 'pdf':
                return DocumentExtractor.extract_pdf(file_content)
            elif ext in ['docx', 'doc']:
                return DocumentExtractor.extract_docx(file_content)
            elif ext in ['xlsx', 'xls']:
                return DocumentExtractor.extract_xlsx(file_content)
            elif ext in ['pptx', 'ppt']:
                return DocumentExtractor.extract_pptx(file_content)
            else:
                # Fallback to UTF-8 text if unknown
                return file_content.decode('utf-8', errors='ignore')
        except Exception as e:
            print(f"[EXTRACTOR ERROR] Failed to extract from {filename}: {e}")
            return ""

    @staticmethod
    def extract_pdf(content: bytes) -> str:
        reader = PyPDF2.PdfReader(io.BytesIO(content))
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text

    @staticmethod
    def extract_docx(content: bytes) -> str:
        doc = docx.Document(io.BytesIO(content))
        return "\n".join([p.text for p in doc.paragraphs])

    @staticmethod
    def extract_xlsx(content: bytes) -> str:
        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        text = ""
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
        return text

    @staticmethod
    def extract_pptx(content: bytes) -> str:
        prs = Presentation(io.BytesIO(content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
